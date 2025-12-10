import os
import json
import base64
from io import BytesIO
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
from openpyxl import load_workbook
import glob
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
import plotly.graph_objects as go
from plotly.io import to_image
import requests
from PIL import Image
import ollama
from typing import List, Dict, Any, Optional
from odf.opendocument import load
from odf.text import P
from odf.table import Table, TableRow, TableCell
import PyPDF2
import zipfile

class OllamaDocumentToPowerPoint:
    def __init__(self, ollama_model: str = "qwen3:30b"):
        """
        Initialize the Ollama-powered PowerPoint generator
        
        Args:
            ollama_model: Ollama model to use (default: qwen3:30b)
        """
        self.presentation = Presentation()
        self.slide_layouts = self.presentation.slide_layouts
        self.ollama_model = ollama_model
        self.ollama_available = self._check_ollama_availability()
        
        print("[FILE] Using standard PowerPoint presentation")
        
        # Comprehensive analysis keywords for universal subject detection
        self.analysis_keywords = {
            # Business & Finance
            'business': ['revenue', 'profit', 'sales', 'cost', 'budget', 'roi', 'growth', 'market', 'strategy', 'investment', 'finance', 'accounting', 'economics', 'pricing', 'margin'],
            
            # Technology & Engineering
            'technical': ['performance', 'efficiency', 'implementation', 'system', 'process', 'software', 'hardware', 'technology', 'engineering', 'development', 'programming', 'algorithm', 'database', 'network'],
            
            # Data & Analytics
            'data': ['average', 'total', 'count', 'percentage', 'trend', 'analysis', 'statistics', 'metrics', 'measurement', 'correlation', 'regression', 'dataset', 'visualization', 'insights'],
            
            # Time & Temporal
            'time': ['month', 'year', 'quarter', 'weekly', 'daily', 'period', 'timeline', 'schedule', 'deadline', 'duration', 'frequency', 'seasonal', 'annual', 'historical'],
            
            # Health & Medical
            'medical': ['patient', 'treatment', 'diagnosis', 'symptoms', 'therapy', 'medication', 'clinical', 'health', 'disease', 'medical', 'healthcare', 'hospital', 'doctor', 'nursing'],
            
            # Education & Research
            'academic': ['research', 'study', 'education', 'learning', 'curriculum', 'student', 'teacher', 'academic', 'university', 'course', 'assessment', 'knowledge', 'methodology', 'theory'],
            
            # Legal & Compliance
            'legal': ['law', 'legal', 'regulation', 'compliance', 'policy', 'contract', 'agreement', 'liability', 'jurisdiction', 'legislation', 'court', 'attorney', 'rights', 'obligations'],
            
            # Science & Research
            'scientific': ['experiment', 'hypothesis', 'research', 'scientific', 'laboratory', 'methodology', 'results', 'conclusion', 'observation', 'theory', 'evidence', 'testing', 'validation', 'peer'],
            
            # Marketing & Communications
            'marketing': ['campaign', 'brand', 'advertising', 'promotion', 'customer', 'audience', 'engagement', 'communication', 'social', 'digital', 'content', 'messaging', 'outreach', 'awareness'],
            
            # Operations & Management
            'operations': ['workflow', 'procedure', 'management', 'operations', 'logistics', 'supply', 'inventory', 'quality', 'control', 'optimization', 'resource', 'planning', 'coordination', 'execution'],
            
            # Human Resources & Personnel
            'hr': ['employee', 'staff', 'personnel', 'hiring', 'training', 'development', 'performance', 'evaluation', 'compensation', 'benefits', 'workplace', 'culture', 'team', 'leadership'],
            
            # Environmental & Sustainability
            'environmental': ['environment', 'sustainability', 'climate', 'energy', 'carbon', 'emission', 'renewable', 'conservation', 'ecological', 'green', 'pollution', 'waste', 'recycling', 'impact'],
            
            # Quality & Standards
            'quality': ['quality', 'standard', 'specification', 'requirement', 'criteria', 'benchmark', 'excellence', 'improvement', 'assurance', 'control', 'testing', 'validation', 'certification', 'audit'],
            
            # Risk & Security
            'risk': ['risk', 'security', 'threat', 'vulnerability', 'safety', 'protection', 'mitigation', 'assessment', 'management', 'contingency', 'emergency', 'prevention', 'monitoring', 'compliance']
        }
        
        # Ollama prompts for different content types
        self.ollama_prompts = {
            'document_summary': """
            Create an executive summary of this document for a PowerPoint presentation. 
            Structure your response as:
            1. Main Purpose (1-2 sentences)
            2. Key Points (3-5 bullet points)
            3. Important Numbers/Data (if any)
            4. Conclusions/Recommendations (2-3 bullet points)
            
            Keep total response under 300 words. Be professional and concise.
            
            Document Content: {content}
            """,
            'detailed_analysis': """
            Provide a detailed analysis of this content for business presentation purposes.
            Focus on:
            - What this document is about
            - Key findings or insights
            - Business implications
            - Actionable recommendations
            
            Format as bullet points. Keep under 400 words.
            
            Content: {content}
            """,
            'excel_insights': """
            Analyze this Excel data and provide business insights for a PowerPoint presentation.
            Focus on:
            - Key trends and patterns
            - Notable numbers or percentages
            - Business implications
            - Recommendations based on data
            
            Format as bullet points. Keep under 300 words.
            
            Data Summary: {content}
            """,
            'slide_title_generator': """
            Create a compelling slide title (max 8 words) for this content.
            Make it professional and descriptive.
            
            Content: {content}
            """,
            'powerbi_analysis': """
            Analyze this Power BI content for presentation purposes.
            Focus on:
            - What insights this dashboard provides
            - Key metrics and KPIs
            - Business value
            - Actionable insights
            
            Format as bullet points. Keep under 300 words.
            
            Power BI Content: {content}
            """,
            'comprehensive_summary': """
            Create a comprehensive executive summary combining insights from multiple documents.
            Structure your response as:
            
            EXECUTIVE SUMMARY
            1. Overall Purpose & Scope (2-3 sentences)
            2. Key Findings (4-6 bullet points)
            3. Critical Data Points (numbers, percentages, trends)
            4. Strategic Recommendations (3-4 bullet points)
            5. Next Steps (2-3 bullet points)
            
            Keep professional tone. Maximum 500 words total.
            
            Combined Content: {content}
            """
        }
    

    
    def _check_ollama_availability(self) -> bool:
        """Check if Ollama is available and the model is accessible"""
        try:
            # Test if Ollama is running and model is available
            response = ollama.list()
            
            # Extract model names from response
            model_names = []
            if hasattr(response, 'models'):
                for model in response.models:
                    if hasattr(model, 'model'):
                        model_names.append(model.model)
            
            print(f"Available Ollama models: {model_names}")
            
            if self.ollama_model in model_names:
                print(f"[OK] Ollama model '{self.ollama_model}' is available")
                return True
            else:
                print(f"[WARNING] Ollama model '{self.ollama_model}' not found.")
                print("Using built-in analysis as fallback")
                return False
        except Exception as e:
            print(f"[WARNING] Ollama not available: {e}")
            print("Using built-in analysis as fallback")
            return False
    
    def call_ollama_api(self, prompt: str, content: str, timeout_seconds: int = 30) -> str:
        """Call Ollama API for content generation with timeout"""
        if not self.ollama_available:
            return self.analyze_content_builtin(content, "summary")
        
        try:
            print(f"[AI] Calling Ollama AI ({self.ollama_model})...")
            full_prompt = prompt.format(content=content[:2000])  # Limit content length
            
            response = ollama.chat(
                model=self.ollama_model,
                messages=[
                    {
                        'role': 'system', 
                        'content': 'You are a business presentation expert. Be very concise (under 200 words).'
                    },
                    {
                        'role': 'user', 
                        'content': full_prompt
                    }
                ],
                options={
                    'num_predict': 200,  # Limit response length
                    'temperature': 0.7,
                    'timeout': timeout_seconds
                }
            )
            
            ai_response = response['message']['content'].strip()
            print(f"[OK] Ollama AI response received ({len(ai_response)} chars)")
            return ai_response
            
        except Exception as e:
            print(f"[WARNING] Ollama API error: {e}")
            print("🔄 Falling back to built-in analysis")
            return self.analyze_content_builtin(content, "summary")

    def analyze_content_builtin(self, content: str, analysis_type: str = "summary") -> str:
        """Built-in content analysis without external API"""
        if not content:
            return "No content to analyze"
        
        content_lower = content.lower()
        
        if analysis_type == "summary":
            return self._create_summary(content)
        elif analysis_type == "insights":
            return self._generate_insights(content_lower)
        elif analysis_type == "title":
            return self._create_title(content)
        else:
            return self._create_summary(content)
    
    def _create_summary(self, content: str) -> str:
        """Create a summary using built-in text analysis"""
        sentences = content.split('.')
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        
        # Take first few sentences and key sentences
        summary_parts = []
        if sentences:
            summary_parts.append(sentences[0])  # First sentence
            
            # Look for sentences with key business words
            for sentence in sentences[1:4]:  # Check next 3 sentences
                sentence_lower = sentence.lower()
                if any(keyword in sentence_lower for keywords in self.analysis_keywords.values() for keyword in keywords):
                    summary_parts.append(sentence)
        
        return ". ".join(summary_parts[:3]) + "." if summary_parts else "Content summary not available"
    
    def _generate_document_summary(self, content: str, content_items: list = None) -> str:
        """Generate a structured summary of the main points from the document"""
        if not content and not content_items:
            return "No content available for summary"
        
        # Use content_items if available (structured content), otherwise parse content
        if content_items:
            return self._create_structured_summary(content_items)
        else:
            return self._create_text_summary(content)
    
    def _create_structured_summary(self, content_items: list) -> str:
        """Create summary from structured content items"""
        summary_points = []
        
        # Group content by type
        paragraphs = []
        tables = []
        lists = []
        headers = []
        
        for item in content_items:
            if item.get('type') == 'paragraph':
                text = item.get('text', '').strip()
                if len(text) > 50:  # Only include substantial paragraphs
                    paragraphs.append(text)
            elif item.get('type') == 'table':
                tables.append(item)
            elif item.get('type') == 'list':
                lists.extend(item.get('items', []))
            elif item.get('type') == 'header':
                headers.append(item.get('text', ''))
        
        # Create summary sections
        if headers:
            summary_points.append("MAIN SECTIONS:")
            for i, header in enumerate(headers[:5], 1):  # Max 5 headers
                summary_points.append(f"{i}. {header}")
        
        if paragraphs:
            summary_points.append("\nKEY POINTS:")
            for i, para in enumerate(paragraphs[:3], 1):  # Max 3 key paragraphs
                # Take first sentence or first 100 characters
                first_sentence = para.split('.')[0] + '.' if '.' in para else para[:100] + '...'
                summary_points.append(f"• {first_sentence}")
        
        if tables:
            summary_points.append(f"\nDATA TABLES: {len(tables)} table(s) with structured information")
        
        if lists:
            summary_points.append("\nLIST ITEMS:")
            for i, item in enumerate(lists[:3], 1):  # Max 3 list items
                summary_points.append(f"• {item}")
        
        return "\n".join(summary_points) if summary_points else "Document structure identified - content ready for presentation"
    
    def _extract_main_points(self, content_items: list) -> str:
        """Extract the main points from content items for presentation"""
        main_points = []
        
        for item in content_items:
            if item.get('type') == 'paragraph':
                text = item.get('text', '').strip()
                if len(text) > 30:  # Only substantial content
                    # Get first sentence or key phrase
                    first_sentence = text.split('.')[0] if '.' in text else text[:100]
                    main_points.append(f"• {first_sentence.strip()}")
        
        # Limit to top 5 points
        return "\n".join(main_points[:5]) if main_points else "• Key content points identified from document"
    
    def _generate_action_items(self, content_items: list) -> str:
        """Generate practical action items based on document content"""
        action_items = []
        
        # Look for action-oriented content
        action_keywords = ['implement', 'create', 'develop', 'establish', 'plan', 'review', 'analyze', 'consider']
        
        for item in content_items:
            text = item.get('text', '').lower()
            for keyword in action_keywords:
                if keyword in text:
                    # Extract sentence containing action keyword
                    sentences = item.get('text', '').split('.')
                    for sentence in sentences:
                        if keyword in sentence.lower():
                            action_items.append(f"• {sentence.strip()}")
                            break
                    break
        
        # If no specific actions found, create generic ones
        if not action_items:
            action_items = [
                "• Review the main points presented",
                "• Identify key priorities for implementation",
                "• Assign responsibilities and timelines",
                "• Schedule follow-up discussions"
            ]
        
        return "\n".join(action_items[:4])  # Max 4 action items
    
    def _create_text_summary(self, content: str) -> str:
        """Create summary from plain text content"""
        if not content:
            return "No text content available"
        
        # Split into paragraphs and sentences
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
        all_sentences = []
        
        for para in paragraphs:
            sentences = [s.strip() for s in para.split('.') if s.strip() and len(s.strip()) > 20]
            all_sentences.extend(sentences)
        
        summary_points = []
        
        # Document overview
        word_count = len(content.split())
        para_count = len(paragraphs)
        summary_points.append(f"DOCUMENT OVERVIEW: {para_count} sections, {word_count} words")
        
        # Key content points
        if all_sentences:
            summary_points.append("\nMAIN POINTS:")
            # Take first sentence from each paragraph (up to 5)
            for i, sentence in enumerate(all_sentences[:5], 1):
                summary_points.append(f"{i}. {sentence}.")
        
        # Content characteristics
        content_lower = content.lower()
        detected_topics = []
        
        for category, keywords in self.analysis_keywords.items():
            if any(word in content_lower for word in keywords):
                detected_topics.append(category.title())
        
        if detected_topics:
            summary_points.append(f"\nTOPIC AREAS: {', '.join(detected_topics[:3])}")
        
        return "\n".join(summary_points)
    
    def _generate_insights(self, content_lower: str) -> str:
        """Generate document summary instead of AI insights - maintained for compatibility"""
        # Redirect to document summary for better presentation focus
        return self._generate_document_summary(content_lower)
    
    def _create_title(self, content: str) -> str:
        """Create a title based on comprehensive content analysis"""
        content_lower = content.lower()
        
        # Smart title generation based on detected categories
        title_patterns = {
            'business': "Business Performance Analysis",
            'technical': "Technical System Review", 
            'data': "Data Analysis Report",
            'medical': "Healthcare Analysis",
            'academic': "Research & Educational Content",
            'legal': "Legal & Compliance Review",
            'scientific': "Scientific Research Analysis",
            'marketing': "Marketing & Communication Strategy",
            'operations': "Operations Management Review",
            'hr': "Human Resources Analysis",
            'environmental': "Environmental Impact Assessment",
            'quality': "Quality Management Review",
            'risk': "Risk & Security Assessment"
        }
        
        # Find the best matching category
        for category, keywords in self.analysis_keywords.items():
            if any(word in content_lower for word in keywords[:5]):  # Check top 5 keywords for efficiency
                return title_patterns.get(category, f"{category.title()} Analysis")
        
        # Fallback: Extract meaningful words from content
        words = content.split()[:8]
        meaningful_words = [w for w in words if len(w) > 3 and w.isalpha() and w.lower() not in ['this', 'that', 'with', 'from', 'they', 'have', 'been', 'will', 'were']]
        
        if meaningful_words:
            return " ".join(meaningful_words[:4]).title() + " Analysis"
        else:
            return "Document Analysis"
    
    def _analyze_excel_data(self, df: pd.DataFrame, sheet_name: str) -> str:
        """Analyze Excel data using Ollama or built-in analysis"""
        # Create comprehensive data summary for Ollama
        data_summary = f"""
        Sheet: {sheet_name}
        Shape: {df.shape[0]} rows, {df.shape[1]} columns
        Columns: {', '.join(df.columns.tolist())}
        """
        
        # Add sample data
        if len(df) > 0:
            data_summary += f"\nSample data (first 3 rows):\n{df.head(3).to_string()}"
        
        # Add data types
        data_summary += f"\nData types:\n{df.dtypes.to_string()}"
        
        # Add numeric summary if available
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            numeric_summary = df[numeric_cols].describe()
            data_summary += f"\nNumeric Summary:\n{numeric_summary.to_string()}"
        
        # Try Ollama first, fallback to built-in analysis
        if self.ollama_available:
            try:
                return self.call_ollama_api(self.ollama_prompts['excel_insights'], data_summary)
            except Exception as e:
                print(f"Ollama Excel analysis error: {e}")
                print("Falling back to built-in analysis")
        
        # Fallback to built-in analysis
        insights = []
        insights.append(f"Sheet '{sheet_name}' contains {df.shape[0]} rows and {df.shape[1]} columns")
        
        if len(numeric_cols) > 0:
            insights.append(f"Found {len(numeric_cols)} numeric columns for analysis")
            for col in numeric_cols[:3]:
                col_data = df[col].dropna()
                if len(col_data) > 0:
                    avg_val = col_data.mean()
                    max_val = col_data.max()
                    min_val = col_data.min()
                    insights.append(f"'{col}': Average {avg_val:.2f}, Range {min_val:.2f} - {max_val:.2f}")
        
        text_cols = df.select_dtypes(include=['object']).columns
        if len(text_cols) > 0:
            insights.append(f"Found {len(text_cols)} text/category columns")
        
        return " • ".join(insights) if insights else "Basic data structure analyzed"
    
    def _analyze_powerbi_builtin(self, metadata: Dict[str, Any]) -> str:
        """Analyze Power BI metadata without external API"""
        analysis_points = []
        
        file_size_mb = metadata.get('file_size', 0) / (1024 * 1024)
        analysis_points.append(f"Power BI file size: {file_size_mb:.1f} MB")
        
        model_files = metadata.get('model_files', [])
        if model_files:
            analysis_points.append(f"Contains {len(model_files)} data model components")
        
        if metadata.get('has_layout'):
            analysis_points.append("Report includes visual layout information")
        
        # Determine likely content based on file size
        if file_size_mb > 50:
            analysis_points.append("Large file suggests complex dashboard with multiple data sources")
        elif file_size_mb > 10:
            analysis_points.append("Medium-sized report likely contains multiple visualizations")
        else:
            analysis_points.append("Compact report with focused data analysis")
        
        analysis_points.append("Recommended: Export key visuals and data for presentation")
        
        return " • ".join(analysis_points)
    
    def extract_odt_content(self, odt_file_path):
        """Extract text content from ODT (OpenOffice) document"""
        try:
            doc = load(odt_file_path)
            content = []
            all_text = []
            
            # Extract paragraphs
            paragraphs = doc.getElementsByType(P)
            for paragraph in paragraphs:
                text = str(paragraph).strip()
                if text and len(text) > 5:  # Filter out very short/empty paragraphs
                    content.append({
                        'type': 'paragraph',
                        'text': text,
                        'style': 'Normal'
                    })
                    all_text.append(text)
            
            # Extract tables
            tables = doc.getElementsByType(Table)
            for table in tables:
                table_data = []
                rows = table.getElementsByType(TableRow)
                for row in rows:
                    cells = row.getElementsByType(TableCell)
                    row_data = []
                    for cell in cells:
                        cell_paragraphs = cell.getElementsByType(P)
                        cell_text = ' '.join([str(p).strip() for p in cell_paragraphs])
                        row_data.append(cell_text)
                    if row_data:  # Only add non-empty rows
                        table_data.append(row_data)
                
                if table_data:
                    content.append({
                        'type': 'table',
                        'data': table_data
                    })
            
            # Add comprehensive document analysis
            if all_text:
                full_text = "\n".join(all_text)
                if self.ollama_available:
                    # Get comprehensive document summary
                    doc_summary = self.call_ollama_api(self.ollama_prompts['document_summary'], full_text)
                    detailed_analysis = self.call_ollama_api(self.ollama_prompts['detailed_analysis'], full_text)
                else:
                    doc_summary = self.analyze_content_builtin(full_text, "summary")
                    detailed_analysis = self.analyze_content_builtin(full_text, "insights")
                
                content.append({
                    'type': 'document_analysis',
                    'executive_summary': doc_summary,
                    'detailed_analysis': detailed_analysis,
                    'word_count': len(full_text.split()),
                    'document_type': 'odt_document'
                })
                
            return content
        except Exception as e:
            print(f"Error reading ODT file {odt_file_path}: {e}")
            return []
    
    def extract_pdf_content(self, pdf_file_path):
        """Extract text content from PDF document with AI analysis"""
        try:
            content = []
            all_text = []
            
            with open(pdf_file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                # Extract text from each page
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    
                    if page_text and page_text.strip():
                        # Clean up the text (remove extra whitespace, fix line breaks)
                        cleaned_text = ' '.join(page_text.split())
                        
                        # Split into paragraphs based on common patterns
                        paragraphs = self._split_pdf_text_into_paragraphs(cleaned_text)
                        
                        for paragraph in paragraphs:
                            if len(paragraph.strip()) > 20:  # Only include substantial paragraphs
                                content.append({
                                    'type': 'paragraph',
                                    'text': paragraph.strip(),
                                    'page': page_num + 1,
                                    'style': 'Normal'
                                })
                                all_text.append(paragraph.strip())
                
                # Add document metadata
                pdf_info = pdf_reader.metadata
                if pdf_info:
                    metadata_text = []
                    if pdf_info.get('/Title'):
                        metadata_text.append(f"Title: {pdf_info.get('/Title')}")
                    if pdf_info.get('/Author'):
                        metadata_text.append(f"Author: {pdf_info.get('/Author')}")
                    if pdf_info.get('/Subject'):
                        metadata_text.append(f"Subject: {pdf_info.get('/Subject')}")
                    
                    if metadata_text:
                        content.append({
                            'type': 'metadata',
                            'text': '\n'.join(metadata_text),
                            'source': 'PDF Properties'
                        })
                
                # Add comprehensive document analysis
                if all_text:
                    full_text = "\n".join(all_text)
                    if self.ollama_available:
                        # Get comprehensive document summary
                        doc_summary = self.call_ollama_api(self.ollama_prompts['document_summary'], full_text)
                        detailed_analysis = self.call_ollama_api(self.ollama_prompts['detailed_analysis'], full_text)
                    else:
                        doc_summary = self.analyze_content_builtin(full_text, "summary")
                        detailed_analysis = self.analyze_content_builtin(full_text, "insights")
                    
                    content.append({
                        'type': 'document_analysis',
                        'executive_summary': doc_summary,
                        'detailed_analysis': detailed_analysis,
                        'word_count': len(full_text.split()),
                        'page_count': num_pages,
                        'document_type': 'pdf_document'
                    })
                    
                return content
                
        except Exception as e:
            print(f"Error reading PDF file {pdf_file_path}: {e}")
            return []
    
    def _split_pdf_text_into_paragraphs(self, text: str) -> List[str]:
        """Split extracted PDF text into meaningful paragraphs"""
        # Common patterns that indicate paragraph breaks in PDFs
        import re
        
        # Replace multiple spaces and normalize whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Split on common paragraph indicators
        # Look for patterns like: ". [Capital letter]", ". \n", etc.
        paragraph_patterns = [
            r'\.\s+([A-Z][a-z])',  # Period followed by capital letter
            r'\.\s*\n\s*([A-Z])',  # Period, newline, capital letter
            r'[.!?]\s{2,}([A-Z])',  # Punctuation with multiple spaces
            r'\n\s*([A-Z][a-z]{2,})'  # Newline followed by word starting with capital
        ]
        
        # Split the text using the patterns
        paragraphs = [text]
        for pattern in paragraph_patterns:
            new_paragraphs = []
            for para in paragraphs:
                splits = re.split(pattern, para)
                if len(splits) > 1:
                    # Rejoin the split parts properly
                    for i in range(0, len(splits) - 1, 2):
                        if i + 1 < len(splits):
                            new_paragraphs.append(splits[i])
                            if i + 2 < len(splits):
                                new_paragraphs.append(splits[i + 1] + splits[i + 2])
                else:
                    new_paragraphs.append(para)
            paragraphs = new_paragraphs
        
        # Filter out very short paragraphs and clean up
        cleaned_paragraphs = []
        for para in paragraphs:
            para = para.strip()
            if len(para) > 50:  # Minimum paragraph length
                cleaned_paragraphs.append(para)
        
        return cleaned_paragraphs if cleaned_paragraphs else [text]
    
    def extract_powerbi_metadata(self, powerbi_file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from Power BI files (.pbix)
        Note: This is a basic implementation. Full .pbix parsing requires specialized tools
        """
        try:
            # Power BI files are ZIP archives, we can extract some basic info
            import zipfile
            
            metadata = {
                'file_name': os.path.basename(powerbi_file_path),
                'file_size': os.path.getsize(powerbi_file_path),
                'type': 'powerbi',
                'extracted_info': []
            }
            
            # Try to extract basic structure (this is simplified)
            with zipfile.ZipFile(powerbi_file_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                
                # Look for data model files
                model_files = [f for f in file_list if 'DataModel' in f or '.json' in f]
                metadata['model_files'] = model_files[:5]  # Limit to first 5
                
                # Try to read layout information if available
                try:
                    if 'Layout' in file_list:
                        layout_data = zip_file.read('Layout').decode('utf-8', errors='ignore')
                        metadata['has_layout'] = True
                        metadata['layout_preview'] = layout_data[:500]  # First 500 chars
                except:
                    metadata['has_layout'] = False
            
            return metadata
            
        except Exception as e:
            print(f"Error reading Power BI file {powerbi_file_path}: {e}")
            return {
                'file_name': os.path.basename(powerbi_file_path),
                'type': 'powerbi',
                'error': str(e),
                'extracted_info': []
            }
    
    def generate_chart_from_data(self, data: pd.DataFrame, chart_type: str = "auto") -> Optional[str]:
        """Generate charts from data and return image path"""
        try:
            # Create a figure
            plt.figure(figsize=(10, 6))
            
            # Auto-detect chart type based on data
            if chart_type == "auto":
                if len(data.columns) >= 2:
                    numeric_cols = data.select_dtypes(include=['number']).columns
                    if len(numeric_cols) >= 1:
                        chart_type = "bar"
                    else:
                        chart_type = "table"
                else:
                    chart_type = "table"
            
            # Generate different chart types
            if chart_type == "bar" and len(data.select_dtypes(include=['number']).columns) > 0:
                numeric_col = data.select_dtypes(include=['number']).columns[0]
                if len(data.columns) > 1:
                    category_col = data.columns[0] if data.columns[0] != numeric_col else data.columns[1]
                    data.plot(x=category_col, y=numeric_col, kind='bar')
                else:
                    data[numeric_col].plot(kind='bar')
                    
            elif chart_type == "line" and len(data.select_dtypes(include=['number']).columns) > 0:
                numeric_cols = data.select_dtypes(include=['number']).columns
                for col in numeric_cols[:3]:  # Max 3 lines
                    plt.plot(data.index, data[col], label=col)
                plt.legend()
                
            elif chart_type == "pie" and len(data.select_dtypes(include=['number']).columns) > 0:
                numeric_col = data.select_dtypes(include=['number']).columns[0]
                plt.pie(data[numeric_col], labels=data.index, autopct='%1.1f%%')
            
            # Save the chart
            chart_path = f"temp_chart_{hash(str(data.values.tolist()))}.png"
            try:
                plt.tight_layout()
            except:
                pass  # Suppress tight layout warnings
            plt.savefig(chart_path, dpi=300, bbox_inches='tight')
            plt.close()
            
            return chart_path
            
        except Exception as e:
            print(f"Error generating chart: {e}")
            plt.close()
            return None
        
    def extract_word_content(self, word_file_path):
        """Extract text content from Word document with AI analysis"""
        try:
            doc = Document(word_file_path)
            content = []
            all_text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append({
                        'type': 'paragraph',
                        'text': paragraph.text.strip(),
                        'style': paragraph.style.name if paragraph.style else 'Normal'
                    })
                    all_text.append(paragraph.text.strip())
            
            # Extract tables if any
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                content.append({
                    'type': 'table',
                    'data': table_data
                })
            
            # Add comprehensive document analysis
            if all_text:
                full_text = "\n".join(all_text)
                if self.ollama_available:
                    # Get comprehensive document summary
                    doc_summary = self.call_ollama_api(self.ollama_prompts['document_summary'], full_text)
                    detailed_analysis = self.call_ollama_api(self.ollama_prompts['detailed_analysis'], full_text)
                else:
                    doc_summary = self.analyze_content_builtin(full_text, "summary")
                    detailed_analysis = self.analyze_content_builtin(full_text, "insights")
                
                content.append({
                    'type': 'document_analysis',
                    'executive_summary': doc_summary,
                    'detailed_analysis': detailed_analysis,
                    'word_count': len(full_text.split()),
                    'document_type': 'word_document'
                })
                
            return content
        except Exception as e:
            print(f"Error reading Word file {word_file_path}: {e}")
            return []
    
    def extract_excel_content(self, excel_file_path, sheet_name=None):
        """Extract data from Excel sheets with AI analysis and chart generation"""
        try:
            # Read with pandas for easy data manipulation
            if sheet_name:
                df = pd.read_excel(excel_file_path, sheet_name=sheet_name)
                sheets_data = {sheet_name: df}
            else:
                sheets_data = pd.read_excel(excel_file_path, sheet_name=None)
            
            content = []
            for sheet, data in sheets_data.items():
                # Convert DataFrame to list of lists for easier handling
                table_data = [data.columns.tolist()]  # Headers
                table_data.extend(data.values.tolist())  # Data rows
                
                # Generate chart if data is suitable
                chart_path = None
                if len(data) > 1 and len(data.select_dtypes(include=['number']).columns) > 0:
                    chart_path = self.generate_chart_from_data(data)
                
                # Generate Ollama insights from data
                ollama_insights = None
                if len(data) > 0:
                    ollama_insights = self._analyze_excel_data(data, sheet)
                
                content.append({
                    'type': 'excel_sheet',
                    'sheet_name': sheet,
                    'data': table_data,
                    'shape': data.shape,
                    'chart_path': chart_path,
                    'ollama_insights': ollama_insights,
                    'dataframe': data  # Keep original DataFrame for further processing
                })
            
            return content
        except Exception as e:
            print(f"Error reading Excel file {excel_file_path}: {e}")
            return []
    
    def extract_pptx_content(self, pptx_file_path):
        """Extract content from existing PowerPoint files with AI enhancement"""
        try:
            source_ppt = Presentation(pptx_file_path)
            content = []
            
            for slide_num, slide in enumerate(source_ppt.slides):
                slide_content = {
                    'type': 'pptx_slide',
                    'slide_number': slide_num + 1,
                    'shapes': []
                }
                
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_content['shapes'].append({
                            'type': 'text',
                            'text': shape.text.strip()
                        })
                        slide_text.append(shape.text.strip())
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # Table
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        slide_content['shapes'].append({
                            'type': 'table',
                            'data': table_data
                        })
                
                # Add Ollama analysis for each slide
                if slide_text:
                    combined_text = "\n".join(slide_text)
                    if self.ollama_available:
                        ollama_title = self.call_ollama_api(self.ollama_prompts['slide_title_generator'], combined_text)
                    else:
                        ollama_title = self.analyze_content_builtin(combined_text, "title")
                    slide_content['ollama_title'] = ollama_title
                
                content.append(slide_content)
            
            return content
        except Exception as e:
            print(f"Error reading PowerPoint file {pptx_file_path}: {e}")
            return []
    
    def process_powerbi_files(self, powerbi_files: List[str]) -> List[Dict[str, Any]]:
        """Process Power BI files and extract metadata"""
        content = []
        
        for powerbi_file in powerbi_files:
            if os.path.exists(powerbi_file):
                print(f"Processing Power BI file: {powerbi_file}")
                metadata = self.extract_powerbi_metadata(powerbi_file)
                
                # Generate Ollama analysis of the Power BI file
                if self.ollama_available:
                    powerbi_info = json.dumps(metadata, indent=2)
                    ollama_analysis = self.call_ollama_api(self.ollama_prompts['powerbi_analysis'], powerbi_info)
                else:
                    ollama_analysis = self._analyze_powerbi_builtin(metadata)
                metadata['ollama_analysis'] = ollama_analysis
                
                content.append(metadata)
        
        return content
    



    

    

    

    

    


    def _estimate_text_height(self, text: str, font_size: int, width_inches: float) -> float:
        """Estimate text height in inches based on content and font size"""
        # Rough estimation: ~65 characters per line for typical fonts
        chars_per_line = int(width_inches * 65 / 9)  # Scale based on slide width
        lines_needed = len(text) / chars_per_line
        line_height = font_size * 1.2 / 72  # Convert points to inches with line spacing
        return lines_needed * line_height
    
    def _fit_content_to_slide(self, content_list: List[str], max_height_inches: float = 4.5) -> List[List[str]]:
        """Split content into multiple slides if needed to fit"""
        slides_content = []
        current_slide_content = []
        current_height = 0
        
        for item in content_list:
            item_text = str(item)
            # Split very long items into smaller chunks
            if len(item_text) > 500:
                chunks = self._split_long_text(item_text, 400)
                for chunk in chunks:
                    estimated_height = self._estimate_text_height(chunk, 16, 8.5)
                    
                    if current_height + estimated_height > max_height_inches and current_slide_content:
                        slides_content.append(current_slide_content)
                        current_slide_content = [chunk]
                        current_height = estimated_height
                    else:
                        current_slide_content.append(chunk)
                        current_height += estimated_height
            else:
                estimated_height = self._estimate_text_height(item_text, 16, 8.5)
                
                if current_height + estimated_height > max_height_inches and current_slide_content:
                    slides_content.append(current_slide_content)
                    current_slide_content = [item_text]
                    current_height = estimated_height
                else:
                    current_slide_content.append(item_text)
                    current_height += estimated_height
        
        if current_slide_content:
            slides_content.append(current_slide_content)
        
        return slides_content if slides_content else [[]]
    
    def _split_long_text(self, text: str, max_length: int = 400) -> List[str]:
        """Split long text into smaller chunks at sentence boundaries"""
        if len(text) <= max_length:
            return [text]
        
        chunks = []
        sentences = text.split('. ')
        current_chunk = ""
        
        for sentence in sentences:
            if len(current_chunk + sentence) <= max_length:
                current_chunk += sentence + ". " if not sentence.endswith('.') else sentence + " "
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence + ". " if not sentence.endswith('.') else sentence + " "
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _adjust_font_size_to_fit(self, text_frame, content_height: float, max_height: float):
        """Dynamically adjust font size to fit content"""
        if content_height <= max_height:
            return 16  # Default size
        
        # Calculate scaling factor
        scale_factor = max_height / content_height
        new_size = max(10, int(16 * scale_factor))  # Minimum 10pt font
        
        # Apply new size to all paragraphs
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                paragraph.font.size = Pt(new_size)
        
        return new_size

    def add_text_slide(self, title, content_list, ai_enhanced=True, slide_type="content"):
        """Add slides with intelligent content fitting"""
        # Convert all content to strings
        string_content = []
        for item in content_list:
            if isinstance(item, dict) and item.get('type') == 'paragraph':
                string_content.append(item['text'])
            elif isinstance(item, str):
                string_content.append(item)
            else:
                string_content.append(str(item))
        
        # Split content across multiple slides if needed
        slides_content = self._fit_content_to_slide(string_content)
        
        for i, slide_content in enumerate(slides_content):
            slide_layout = self.slide_layouts[1]  # Title and Content layout
            slide = self.presentation.slides.add_slide(slide_layout)
        
            # Create slide title
            slide_title = title
            if len(slides_content) > 1:
                slide_title += f" ({i+1}/{len(slides_content)})"
                
            # Enhance title with Ollama analysis if available (only for first slide)
            if ai_enhanced and i == 0 and slide_content:
                content_preview = str(slide_content[:2])
                if self.ollama_available:
                    enhanced_title = self.call_ollama_api(self.ollama_prompts['slide_title_generator'], content_preview)
                    if enhanced_title and len(enhanced_title) < 80:
                        slide_title = enhanced_title
                        if len(slides_content) > 1:
                            slide_title += f" ({i+1}/{len(slides_content)})"
            
            # Set title styling
            title_shape = slide.shapes.title
            title_shape.text = slide_title
            
            # Style title
            title_paragraph = title_shape.text_frame.paragraphs[0]
            title_paragraph.font.size = Pt(32 if slide_type == "title" else 24)
            title_paragraph.font.bold = True
            
            # Add content with smart fitting
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = None  # Disable auto-sizing to control manually
            
            # Add content to text frame
            total_content_length = sum(len(str(item)) for item in slide_content)
            
            for j, item in enumerate(slide_content):
                item_text = str(item).strip()
                if not item_text:
                    continue
                    
                if j == 0:
                    tf.text = item_text
                else:
                    p = tf.add_paragraph()
                    p.text = item_text
                    p.level = 0
            
            # Apply smart font sizing
            content_height = self._estimate_text_height(
                "\n".join(str(item) for item in slide_content), 
                16, 8.5
            )
            
            font_size = 16
            if content_height > 4.5:  # If content is too tall
                font_size = self._adjust_font_size_to_fit(tf, content_height, 4.5)
            
            # Apply font styling to all paragraphs
            for paragraph in tf.paragraphs:
                if paragraph.text.strip():
                    paragraph.font.size = Pt(font_size)
                    
            print(f"  [FILE] Created slide: {slide_title} (Font: {font_size}pt, Items: {len(slide_content)})")
    
    def _create_condensed_summary(self, content_list: List[str], max_items: int = 8) -> List[str]:
        """Create a condensed version of content for dense slides"""
        if len(content_list) <= max_items:
            return content_list
        
        # Group similar items and create summaries
        condensed = []
        
        # Take most important items (first few and last few)
        important_items = content_list[:max_items//2]
        if len(content_list) > max_items:
            important_items.extend(content_list[-(max_items//2):])
            
        for item in important_items:
            # Truncate very long items
            item_text = str(item)
            if len(item_text) > 200:
                item_text = item_text[:197] + "..."
            condensed.append(item_text)
        
        # Add summary note if content was condensed
        if len(content_list) > max_items:
            condensed.append(f"... and {len(content_list) - max_items} additional items")
        
        return condensed

    def add_overflow_content_slide(self, title: str, content_items: List[str], slide_number: int = 1):
        """Add a continuation slide for overflow content"""
        slide_layout = self.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Create continuation title
        continuation_title = f"{title} - Additional Details ({slide_number})"
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = continuation_title
        
        # Style title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(22)  # Slightly smaller for continuation
        title_paragraph.font.bold = True
        
        # Add content with better spacing
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        tf.word_wrap = True
        
        # Condense content for continuation slides
        condensed_items = self._create_condensed_summary(content_items, max_items=10)
        
        for i, item in enumerate(condensed_items):
            item_text = str(item).strip()
            if not item_text:
                continue
                
            if i == 0:
                tf.text = f"• {item_text}"
            else:
                p = tf.add_paragraph()
                p.text = f"• {item_text}"
                p.level = 0
        
        # Style content with smaller font for continuation
        for paragraph in tf.paragraphs:
            if paragraph.text.strip():
                paragraph.font.size = Pt(14)  # Smaller for more content
        
        print(f"  [FILE] Added continuation slide: {continuation_title} ({len(condensed_items)} items)")

    def add_executive_summary_slide(self, title: str, summary_data: Dict[str, Any]):
        """Add a comprehensive executive summary slide with content fitting"""
        slide_layout = self.slide_layouts[5]  # Blank layout for custom design
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Executive Summary Box
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        
        # Prepare and fit content
        summary_content = ""
        if summary_data.get('executive_summary'):
            summary_content = summary_data['executive_summary']
            
            # Truncate if too long (estimate ~2000 chars for good fit)
            if len(summary_content) > 2000:
                summary_content = summary_content[:1997] + "..."
                print(f"  [WARNING] Executive summary truncated to fit slide")
        
        # Add content with header
        if summary_content:
            summary_frame.text = "EXECUTIVE SUMMARY\n\n" + summary_content
        else:
            summary_frame.text = "EXECUTIVE SUMMARY\n\nSummary content not available."
        
        # Calculate optimal font size based on content length
        content_length = len(summary_frame.text)
        if content_length > 1500:
            body_font_size = 12
            header_font_size = 16
        elif content_length > 1000:
            body_font_size = 13
            header_font_size = 17
        else:
            body_font_size = 14
            header_font_size = 18
        
        # Style the summary with adaptive sizing
        for paragraph in summary_frame.paragraphs:
            if "EXECUTIVE SUMMARY" in paragraph.text:
                paragraph.font.size = Pt(header_font_size)
                paragraph.font.bold = True
            else:
                paragraph.font.size = Pt(body_font_size)
        
        # Add document stats in corner
        if summary_data.get('word_count'):
            stats_box = slide.shapes.add_textbox(Inches(7), Inches(6.5), Inches(2), Inches(0.8))
            stats_frame = stats_box.text_frame
            stats_frame.text = f"Document Stats:\nWords: {summary_data['word_count']}\nType: {summary_data.get('document_type', 'Document')}"
            for para in stats_frame.paragraphs:
                para.font.size = Pt(10)
    
    def add_document_summary_slide(self, title: str, summary_data: Dict[str, Any]):
        """Add a document summary slide focused on main points and overview"""
        slide_layout = self.slide_layouts[5]  # Blank layout for custom design
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Document Summary Box
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        
        # Prepare summary content
        summary_content = ""
        if summary_data.get('summary'):
            summary_content = summary_data['summary']
            
            # Truncate if too long for readability
            if len(summary_content) > 1800:
                summary_content = summary_content[:1797] + "..."
                print(f"  [WARNING] Document summary truncated to fit slide")
        
        # Add document overview information
        doc_info = []
        if summary_data.get('word_count'):
            doc_info.append(f"Word Count: {summary_data['word_count']}")
        if summary_data.get('section_count'):
            doc_info.append(f"Sections: {summary_data['section_count']}")
        
        # Build final content
        final_content = "DOCUMENT OVERVIEW\n\n"
        if doc_info:
            final_content += " | ".join(doc_info) + "\n\n"
        
        if summary_content:
            final_content += summary_content
        else:
            final_content += "Document summary ready for presentation"
        
        summary_frame.text = final_content
        
        # Set font formatting
        content_length = len(final_content)
        if content_length > 1500:
            body_font_size = 12
            header_font_size = 16
        elif content_length > 1000:
            body_font_size = 13
            header_font_size = 17
        else:
            body_font_size = 14
            header_font_size = 18
        
        # Format paragraphs
        for i, para in enumerate(summary_frame.paragraphs):
            if i == 0:  # Header
                para.font.size = Pt(header_font_size)
                para.font.bold = True
            else:
                para.font.size = Pt(body_font_size)
        
        print(f"  [INFO] Added document summary slide with {len(final_content)} characters")
    
    def _calculate_optimal_chart_size(self, slide_width: float = 10.0, slide_height: float = 7.5, 
                                     has_insights: bool = False) -> dict:
        """Calculate optimal chart dimensions to prevent overlapping"""
        # Reserve space for title (1 inch from top)
        available_height = slide_height - 1.5  # Leave space for title and margins
        available_width = slide_width - 1.0    # Leave margins
        
        if has_insights:
            # Split slide: chart on left, insights on right
            chart_width = available_width * 0.6   # 60% for chart
            insights_width = available_width * 0.35  # 35% for insights, 5% gap
            chart_left = 0.5
            insights_left = chart_left + chart_width + 0.25  # Small gap
        else:
            # Full width for chart
            chart_width = available_width
            insights_width = 0
            chart_left = 0.5
            insights_left = 0
        
        # Maintain aspect ratio but fit within bounds
        max_chart_height = available_height * 0.85  # Leave some bottom margin
        
        return {
            'chart': {
                'left': chart_left,
                'top': 1.2,
                'width': chart_width,
                'height': max_chart_height
            },
            'insights': {
                'left': insights_left,
                'top': 1.2, 
                'width': insights_width,
                'height': max_chart_height
            }
        }
    
    def _fit_text_to_box(self, text_frame, max_height_inches: float):
        """Dynamically adjust text size to fit within box bounds"""
        if not text_frame.text.strip():
            return 12
            
        # Calculate text density
        total_chars = len(text_frame.text)
        
        # Determine optimal font size based on content and space
        if total_chars > 800:
            font_size = 10
        elif total_chars > 400:
            font_size = 11
        else:
            font_size = 12
            
        # Apply font size to all paragraphs
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                paragraph.font.size = Pt(font_size)
                
        return font_size

    def add_chart_slide(self, title: str, chart_path: str, insights: str = None):
        """Add a slide with perfectly fitted chart and no overlapping text"""
        slide_layout = self.slide_layouts[6]  # Blank layout for precise control
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Calculate optimal dimensions
        dimensions = self._calculate_optimal_chart_size(has_insights=bool(insights))
        
        # Add title with standard styling
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        
        try:
            # Add chart with calculated dimensions
            if chart_path and os.path.exists(chart_path):
                chart_dims = dimensions['chart']
                slide.shapes.add_picture(
                    chart_path, 
                    Inches(chart_dims['left']), 
                    Inches(chart_dims['top']), 
                    Inches(chart_dims['width']), 
                    Inches(chart_dims['height'])
                )
                
                # Add insights text with perfect fitting
                if insights:
                    insights_dims = dimensions['insights']
                    insights_box = slide.shapes.add_textbox(
                        Inches(insights_dims['left']), 
                        Inches(insights_dims['top']), 
                        Inches(insights_dims['width']), 
                        Inches(insights_dims['height'])
                    )
                    insights_frame = insights_box.text_frame
                    insights_frame.word_wrap = True
                    
                    # Smart text condensation for long insights
                    if len(insights) > 600:
                        # Truncate very long insights to fit
                        insights = insights[:597] + "..."
                        
                    insights_frame.text = f"Key Insights:\n\n{insights}"
                    
                    # Apply optimal font sizing
                    font_size = self._fit_text_to_box(insights_frame, insights_dims['height'])
                    print(f"  [CHART] Chart slide created - Chart: {chart_dims['width']:.1f}\"×{chart_dims['height']:.1f}\", Insights: {font_size}pt")
                else:
                    print(f"  [CHART] Full-width chart slide created - {chart_dims['width']:.1f}\"×{chart_dims['height']:.1f}\"")
            else:
                # Chart file doesn't exist, create text-only slide
                self._create_chart_fallback_slide(slide, title, insights, chart_path)
                    
        except Exception as e:
            # Robust fallback with proper spacing
            self._create_chart_fallback_slide(slide, title, insights, chart_path, str(e))
    
    def _create_chart_fallback_slide(self, slide, title: str, insights: str, chart_path: str, error: str = None):
        """Create fallback slide when chart cannot be loaded"""
        error_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        error_frame = error_box.text_frame
        error_frame.word_wrap = True
        
        error_text = f"[CHART] Chart Analysis: {title}\n\n"
        if error:
            error_text += f"Chart could not be loaded: {error}\n\n"
        else:
            error_text += f"Chart file not found: {chart_path}\n\n"
            
        if insights:
            error_text += f"Analysis:\n{insights[:400]}{'...' if len(insights) > 400 else ''}"
        else:
            error_text += "No additional analysis available."
            
        error_frame.text = error_text
        self._fit_text_to_box(error_frame, 4.0)
        print(f"  [WARNING] Chart slide created with text fallback")
    
    def add_powerbi_slide(self, powerbi_metadata: Dict[str, Any]):
        """Add a slide for Power BI file information"""
        slide_layout = self.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        slide.shapes.title.text = f"Power BI Analysis: {powerbi_metadata.get('file_name', 'Unknown')}"
        
        # Add content
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        
        # Add basic info
        info_items = [
            f"File Size: {powerbi_metadata.get('file_size', 0) / 1024 / 1024:.1f} MB",
            f"Model Files Found: {len(powerbi_metadata.get('model_files', []))}",
            f"Has Layout Information: {'Yes' if powerbi_metadata.get('has_layout') else 'No'}"
        ]
        
        for item in info_items:
            p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text else tf.add_paragraph()
            p.text = item
            p.level = 0
        
        # Add Ollama analysis if available
        if powerbi_metadata.get('ollama_analysis'):
            analysis_p = tf.add_paragraph()
            analysis_p.text = f"Analysis: {powerbi_metadata['ollama_analysis']}"
            analysis_p.level = 0
    
    def add_table_slide(self, title, table_data):
        """Add a slide with a table"""
        slide_layout = self.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Add table
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), 
                                                Inches(9), Inches(5))
            table = table_shape.table
            
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    if j < len(table.rows[i].cells):
                        table.rows[i].cells[j].text = str(cell_data) if cell_data is not None else ""
    
    def process_files(self, input_folder=None, word_files=None, excel_files=None, pptx_files=None, powerbi_files=None, odt_files=None, pdf_files=None):
        """Process multiple files and create presentation with built-in intelligence"""
        all_content = []
        
        # If input folder is provided, scan for files
        if input_folder and os.path.exists(input_folder):
            # Get all files and filter out temporary files (starting with ~$)
            word_files = word_files or [f for f in glob.glob(os.path.join(input_folder, "*.docx")) if not os.path.basename(f).startswith("~$")]
            excel_files = excel_files or [f for f in glob.glob(os.path.join(input_folder, "*.xlsx")) if not os.path.basename(f).startswith("~$")]
            pptx_files = pptx_files or [f for f in glob.glob(os.path.join(input_folder, "*.pptx")) if not os.path.basename(f).startswith("~$")]
            powerbi_files = powerbi_files or [f for f in glob.glob(os.path.join(input_folder, "*.pbix")) if not os.path.basename(f).startswith("~$")]
            odt_files = odt_files or [f for f in glob.glob(os.path.join(input_folder, "*.odt")) if not os.path.basename(f).startswith("~$")]
            pdf_files = pdf_files or [f for f in glob.glob(os.path.join(input_folder, "*.pdf")) if not os.path.basename(f).startswith("~$")]
        
        # Process Word files
        if word_files:
            for word_file in word_files:
                if os.path.exists(word_file):
                    print(f"Processing Word file: {word_file}")
                    content = self.extract_word_content(word_file)
                    all_content.append({
                        'source': os.path.basename(word_file),
                        'type': 'word',
                        'content': content
                    })
        
        # Process Excel files
        if excel_files:
            for excel_file in excel_files:
                if os.path.exists(excel_file):
                    print(f"Processing Excel file: {excel_file}")
                    content = self.extract_excel_content(excel_file)
                    all_content.append({
                        'source': os.path.basename(excel_file),
                        'type': 'excel',
                        'content': content
                    })
        
        # Process PowerPoint files
        if pptx_files:
            for pptx_file in pptx_files:
                if os.path.exists(pptx_file):
                    print(f"Processing PowerPoint file: {pptx_file}")
                    content = self.extract_pptx_content(pptx_file)
                    all_content.append({
                        'source': os.path.basename(pptx_file),
                        'type': 'pptx',
                        'content': content
                    })
        
        # Process Power BI files
        if powerbi_files:
            powerbi_content = self.process_powerbi_files(powerbi_files)
            for pb_content in powerbi_content:
                all_content.append({
                    'source': pb_content.get('file_name', 'Unknown'),
                    'type': 'powerbi',
                    'content': [pb_content]
                })
        
        # Process ODT files
        if odt_files:
            for odt_file in odt_files:
                if os.path.exists(odt_file):
                    print(f"Processing ODT file: {odt_file}")
                    content = self.extract_odt_content(odt_file)
                    all_content.append({
                        'source': os.path.basename(odt_file),
                        'type': 'odt',
                        'content': content
                    })
        
        # Process PDF files
        if pdf_files:
            for pdf_file in pdf_files:
                if os.path.exists(pdf_file):
                    print(f"Processing PDF file: {pdf_file}")
                    content = self.extract_pdf_content(pdf_file)
                    all_content.append({
                        'source': os.path.basename(pdf_file),
                        'type': 'pdf',
                        'content': content
                    })
        
        return all_content
    
    def process_custom_text(self, user_text: str, text_title: str = "Custom Ideas", enable_ai_analysis: bool = False):
        """Process user's custom text input and create presentation content"""
        if not user_text or not user_text.strip():
            print("[WARNING] No text provided for processing")
            return []
        
        print(f"Processing custom text input: '{text_title}'")
        content = []
        
        # Split text into meaningful sections
        sections = self._split_custom_text_into_sections(user_text)
        
        # Process each section
        for i, section in enumerate(sections):
            if len(section.strip()) < 10:  # Skip very short sections
                continue
                
            content.append({
                'type': 'paragraph',
                'text': section.strip(),
                'section': i + 1,
                'style': 'Normal'
            })
        
        # Add document summary instead of AI insights
        if enable_ai_analysis and content:  # Keep parameter for compatibility but change behavior
            full_text = "\n".join([item['text'] for item in content])
            
            # Generate structured document summary
            doc_summary = self._generate_document_summary(full_text, content)
            
            # Create main points overview
            main_points = self._extract_main_points(content)
            
            # Generate simple action items based on content structure
            action_items = self._generate_action_items(content)
            
            content.append({
                'type': 'document_summary',
                'summary': doc_summary,
                'main_points': main_points,
                'action_items': action_items,
                'word_count': len(full_text.split()),
                'section_count': len(sections),
                'document_type': 'custom_text'
            })
        
        return [{
            'source': text_title,
            'type': 'custom_text',
            'content': content
        }]
    
    def _split_custom_text_into_sections(self, text: str) -> List[str]:
        """Split user text into logical sections for better presentation"""
        import re
        
        # Clean up the text
        text = text.strip()
        
        # Try different splitting strategies
        sections = []
        
        # Strategy 1: Split by numbered lists (1., 2., etc.)
        numbered_pattern = r'(\d+\.\s*)'
        if re.search(numbered_pattern, text):
            parts = re.split(numbered_pattern, text)
            current_section = ""
            for i, part in enumerate(parts):
                if re.match(r'\d+\.\s*', part):
                    if current_section.strip():
                        sections.append(current_section.strip())
                    current_section = part
                else:
                    current_section += part
            if current_section.strip():
                sections.append(current_section.strip())
        
        # Strategy 2: Split by bullet points (•, -, *, etc.)
        elif re.search(r'[•\-\*]\s*', text):
            bullet_pattern = r'([•\-\*]\s*)'
            parts = re.split(bullet_pattern, text)
            current_section = ""
            for part in parts:
                if re.match(r'[•\-\*]\s*', part):
                    if current_section.strip():
                        sections.append(current_section.strip())
                    current_section = part
                else:
                    current_section += part
            if current_section.strip():
                sections.append(current_section.strip())
        
        # Strategy 3: Split by double line breaks
        elif '\n\n' in text:
            sections = [section.strip() for section in text.split('\n\n') if section.strip()]
        
        # Strategy 4: Split by single line breaks (if text has clear line structure)
        elif text.count('\n') > 2:
            lines = text.split('\n')
            sections = [line.strip() for line in lines if line.strip() and len(line.strip()) > 10]
        
        # Strategy 5: Split by sentences for longer single paragraphs
        else:
            sentences = re.split(r'[.!?]+\s+', text)
            # Group sentences into sections of 2-3 sentences each
            current_section = ""
            sentence_count = 0
            for sentence in sentences:
                if sentence.strip():
                    current_section += sentence.strip() + ". "
                    sentence_count += 1
                    if sentence_count >= 2 and len(current_section) > 100:
                        sections.append(current_section.strip())
                        current_section = ""
                        sentence_count = 0
            if current_section.strip():
                sections.append(current_section.strip())
        
        # Fallback: return original text as single section
        if not sections or all(len(s) < 10 for s in sections):
            sections = [text]
        
        return sections
    
    def create_presentation_from_text(self, user_text: str, presentation_title: str = "My Ideas Presentation", text_title: str = "Custom Ideas", output_file: str = None, enable_ai_analysis: bool = False):
        """Create a complete presentation from user's custom text input"""
        if not output_file:
            output_file = f"{presentation_title.lower().replace(' ', '_')}.pptx"
        
        # Process the custom text
        all_content = self.process_custom_text(user_text, text_title, enable_ai_analysis)
        
        if not all_content:
            print("[ERROR] No content generated from input text")
            return None
        
        # Create the presentation
        return self.create_presentation(all_content, output_file, presentation_title)
    
    def create_presentation(self, all_content, output_file="ai_combined_presentation.pptx", presentation_title="AI-Enhanced Document Analysis"):
        """Create the final AI-enhanced PowerPoint presentation"""
        # Add title slide
        title_slide_layout = self.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = presentation_title
        
        # Style title slide
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        
        # Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"Comprehensive Document Analysis with AI Insights"
        subtitle_para = subtitle_shape.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        
        # Add overview slide
        overview_content = []
        for source_data in all_content:
            overview_content.append(f"• {source_data['source']} ({source_data['type'].upper()})")
        
        self.add_text_slide("Document Overview", overview_content, ai_enhanced=False)
        
        # Process each source file
        for source_data in all_content:
            source_name = source_data['source']
            content_type = source_data['type']
            content = source_data['content']
            
            # Add section slide for each source file
            self.add_text_slide(f"Analysis: {source_name}", [f"Source type: {content_type.upper()}", "Intelligent analysis included"], ai_enhanced=False)
            
            # Process content based on type
            if content_type == 'word':
                self._process_word_content(content, source_name)
            elif content_type == 'excel':
                self._process_excel_content(content, source_name)
            elif content_type == 'pptx':
                self._process_pptx_content(content, source_name)
            elif content_type == 'powerbi':
                self._process_powerbi_content(content, source_name)
            elif content_type == 'odt':
                self._process_odt_content(content, source_name)
            elif content_type == 'pdf':
                self._process_pdf_content(content, source_name)
            elif content_type == 'custom_text':
                self._process_custom_text_content(content, source_name)
        
        # Add comprehensive summary slide
        self._add_comprehensive_summary_slide(all_content)
        
        # Clean up temporary chart files
        self._cleanup_temp_files()
        
        # Save presentation
        self.presentation.save(output_file)
        print(f"Intelligent presentation saved as: {output_file}")
        return output_file
    
    def _process_word_content(self, content, source_name):
        """Process Word document content into slides with AI enhancement"""
        text_items = []
        
        for item in content:
            if item['type'] == 'paragraph':
                text_items.append(item['text'])
                # Create slide every 5 paragraphs or when we hit a heading
                if len(text_items) >= 5 or 'Heading' in item.get('style', ''):
                    if text_items:
                        self.add_text_slide(f"{source_name} - Content", text_items)
                        text_items = []
            
            elif item['type'] == 'table':
                # Add remaining text items first
                if text_items:
                    self.add_text_slide(f"{source_name} - Content", text_items)
                    text_items = []
                # Add table slide
                self.add_table_slide(f"{source_name} - Table", item['data'])
            
            elif item['type'] == 'document_analysis':
                # Add comprehensive document analysis slides
                # Executive Summary slide
                self.add_executive_summary_slide(f"{source_name} - Executive Summary", item)
                
                # Detailed Analysis slide
                if item.get('detailed_analysis'):
                    self.add_text_slide(f"{source_name} - Detailed Analysis", [item['detailed_analysis']], ai_enhanced=False)
        
        # Add any remaining text items
        if text_items:
            self.add_text_slide(f"{source_name} - Content", text_items)
    
    def _process_excel_content(self, content, source_name):
        """Process Excel content into slides with charts and AI insights"""
        for item in content:
            if item['type'] == 'excel_sheet':
                sheet_name = item['sheet_name']
                table_data = item['data']
                chart_path = item.get('chart_path')
                ollama_insights = item.get('ollama_insights')
                
                # Add chart slide if chart was generated
                if chart_path:
                    chart_title = f"{source_name} - {sheet_name} Visualization"
                    self.add_chart_slide(chart_title, chart_path, ollama_insights)
                
                # Create slide with sheet data (limit rows for readability)
                display_data = table_data[:11] if len(table_data) > 11 else table_data  # Header + 10 rows max
                title = f"{source_name} - {sheet_name} Data"
                self.add_table_slide(title, display_data)
                
                # Add insights slide if available
                if ollama_insights and not chart_path:  # Only if not already shown in chart slide
                    self.add_text_slide(f"{source_name} - {sheet_name} AI Insights", [ollama_insights], ai_enhanced=False)
    
    def _process_pptx_content(self, content, source_name):
        """Process PowerPoint content into slides with AI enhancement"""
        for slide_data in content:
            if slide_data['type'] == 'pptx_slide':
                slide_num = slide_data['slide_number']
                shapes = slide_data['shapes']
                ollama_title = slide_data.get('ollama_title')
                
                text_content = []
                tables = []
                
                for shape in shapes:
                    if shape['type'] == 'text':
                        text_content.append(shape['text'])
                    elif shape['type'] == 'table':
                        tables.append(shape['data'])
                
                # Add text content with AI-enhanced title if available
                if text_content:
                    title = f"{source_name} - Slide {slide_num}"
                    if ollama_title:
                        title += f": {ollama_title}"
                    self.add_text_slide(title, text_content, ai_enhanced=False)
                
                # Add tables
                for i, table in enumerate(tables):
                    self.add_table_slide(f"{source_name} - Slide {slide_num} Table {i+1}", table)
    
    def _process_powerbi_content(self, content, source_name):
        """Process Power BI content into slides"""
        for powerbi_data in content:
            self.add_powerbi_slide(powerbi_data)
    
    def _process_odt_content(self, content, source_name):
        """Process ODT document content into slides (same as Word processing)"""
        # ODT files have the same structure as Word files after extraction
        self._process_word_content(content, source_name)
    
    def _process_pdf_content(self, content, source_name):
        """Process PDF document content into slides with AI enhancement"""
        text_items = []
        metadata_items = []
        
        for item in content:
            if item['type'] == 'paragraph':
                text_items.append(f"[Page {item.get('page', '?')}] {item['text']}")
                # Create slide every 5 paragraphs to manage content length
                if len(text_items) >= 5:
                    if text_items:
                        self.add_text_slide(f"{source_name} - Content", text_items)
                        text_items = []
            
            elif item['type'] == 'metadata':
                metadata_items.append(item['text'])
            
            elif item['type'] == 'document_analysis':
                # Add remaining text items first
                if text_items:
                    self.add_text_slide(f"{source_name} - Content", text_items)
                    text_items = []
                
                # Add PDF metadata slide if available
                if metadata_items:
                    self.add_text_slide(f"{source_name} - Document Info", metadata_items)
                
                # Add comprehensive document analysis slides
                # Executive Summary slide
                self.add_executive_summary_slide(f"{source_name} - Executive Summary", item)
                
                # Detailed Analysis slide
                if item.get('detailed_analysis'):
                    analysis_info = [item['detailed_analysis']]
                    if item.get('page_count'):
                        analysis_info.append(f"Document contains {item['page_count']} pages")
                    self.add_text_slide(f"{source_name} - Detailed Analysis", analysis_info, ai_enhanced=False)
        
        # Add any remaining text items
        if text_items:
            self.add_text_slide(f"{source_name} - Content", text_items)
    
    def _process_custom_text_content(self, content, source_name):
        """Process custom user text input into slides with AI enhancement"""
        text_items = []
        
        for item in content:
            if item['type'] == 'paragraph':
                section_num = item.get('section', 1)
                text_items.append(f"[Section {section_num}] {item['text']}")
                
                # Create slide every 3 sections to keep slides focused
                if len(text_items) >= 3:
                    if text_items:
                        self.add_text_slide(f"{source_name} - Ideas", text_items)
                        text_items = []
            
            elif item['type'] == 'document_summary':
                # Add remaining text items first
                if text_items:
                    self.add_text_slide(f"{source_name} - Content", text_items)
                    text_items = []
                
                # Add document summary slide
                self.add_document_summary_slide(f"{source_name} - Summary", item)
                
                # Add main points slide
                if item.get('main_points'):
                    main_points_info = [item['main_points']]
                    if item.get('word_count'):
                        main_points_info.append(f"Document contains {item['word_count']} words across {item.get('section_count', 'multiple')} sections")
                    self.add_text_slide(f"{source_name} - Key Points", main_points_info, ai_enhanced=False)
                
                # Add action items slide if available
                if item.get('action_items'):
                    action_info = [item['action_items']]
                    action_info.append("\nNext Steps:")
                    action_info.append("• Review and prioritize action items")
                    action_info.append("• Develop implementation timeline")
                    action_info.append("• Identify required resources")
                    self.add_text_slide(f"{source_name} - Action Items", action_info, ai_enhanced=False)
        
        # Add any remaining text items
        if text_items:
            self.add_text_slide(f"{source_name} - Ideas", text_items)
    
    def _add_comprehensive_summary_slide(self, all_content):
        """Add a comprehensive summary slide with AI-enhanced insights"""
        # Gather all analysis content
        all_summaries = []
        all_insights = []
        
        total_files = len(all_content)
        file_types = {}
        
        for source_data in all_content:
            file_type = source_data['type']
            if file_type in file_types:
                file_types[file_type] += 1
            else:
                file_types[file_type] = 1
            
            # Collect analysis data
            for content_item in source_data.get('content', []):
                if content_item.get('type') == 'document_analysis':
                    if content_item.get('executive_summary'):
                        all_summaries.append(content_item['executive_summary'])
                    if content_item.get('detailed_analysis'):
                        all_insights.append(content_item['detailed_analysis'])
        
        # Create comprehensive summary
        combined_content = "\n".join(all_summaries + all_insights)
        
        if self.ollama_available and combined_content:
            # Use AI to create comprehensive summary
            comprehensive_summary = self.call_ollama_api(
                self.ollama_prompts['comprehensive_summary'], 
                combined_content
            )
            self.add_text_slide("Comprehensive Executive Summary", [comprehensive_summary], ai_enhanced=False)
        else:
            # Fallback summary
            summary_content = []
            summary_content.append(f"Document Analysis Summary - {total_files} files processed")
            summary_content.append("")
            
            for file_type, count in file_types.items():
                summary_content.append(f"• {count} {file_type.upper()} file(s) analyzed")
            
            summary_content.append("")
            summary_content.append("Key Deliverables:")
            summary_content.append("• Executive summaries for all documents")
            summary_content.append("• Detailed analysis and insights")
            summary_content.append("• Data visualizations where applicable")
            summary_content.append("• Professional presentation")
            
            self.add_text_slide("Executive Summary", summary_content, ai_enhanced=False)
    
    def _cleanup_temp_files(self):
        """Clean up temporary chart files"""
        temp_files = glob.glob("temp_chart_*.png")
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass  # Ignore errors in cleanup


# Example usage - using files from wrapper test folder
def main():
    # Create converter instance with Ollama
    converter = OllamaDocumentToPowerPoint(
        ollama_model="qwen3:30b"
    )
    
    # Automatically process all files from wrapper test folder
    wrapper_test_folder = r"c:\Users\aidan\OneDrive\Desktop\wrapper test"
    
    print(f"[SEARCH] Scanning folder: {wrapper_test_folder}")
    
    # Method 1: Process all files in the wrapper test folder
    all_content = converter.process_files(input_folder=wrapper_test_folder)
    
    if not all_content:
        print("[ERROR] No compatible files found in the folder!")
        print("Looking for: .docx, .xlsx, .pptx, .pbix, .odt files")
        return
    
    print(f"📁 Found {len(all_content)} file(s) to process")
    
    # Create the comprehensive Ollama-powered presentation
    output_file = converter.create_presentation(
        all_content, 
        "Aidan_Gonzales_TP11_AI_Summary.pptx",
        "Aidan Gonzales TP11 - Comprehensive Document Analysis"
    )
    
    print(f"[OK] Ollama-powered presentation created: {output_file}")


# Alternative method - specify files manually if needed
def main_manual():
    converter = OllamaDocumentToPowerPoint(
        ollama_model="qwen3:30b"
    )
    
    # Manually specify the files from wrapper test folder
    base_path = r"c:\Users\aidan\OneDrive\Desktop\wrapper test"
    
    odt_files = [
        os.path.join(base_path, "Aidan_Gonzales TP11.odt")
    ]
    excel_files = [
        os.path.join(base_path, "Aidan_Gonzales TP11.xlsx")
    ]
    
    # Process the specific files
    all_content = converter.process_files(
        odt_files=odt_files,
        excel_files=excel_files
    )
    
    # Create the presentation
    output_file = converter.create_presentation(
        all_content, 
        "Aidan_Gonzales_TP11_Ollama_Manual.pptx",
        "Aidan Gonzales TP11 - Ollama Manual Processing"
    )
    
    print(f"[OK] Ollama manual presentation created: {output_file}")



# Text Input Demo Function
def create_presentation_from_custom_text():
    """Demo function showing how to create presentations from custom text input"""
    print("=== Custom Text Input Demo ===")
    
    # Example custom text (in real usage, this would come from user input)
    sample_text = """
    Project Innovation Ideas:
    
    1. Implement AI-powered customer service chatbot
    - Reduce response time by 80%
    - Available 24/7 for customer support
    - Learn from interactions to improve over time
    
    2. Develop mobile application for remote work
    - Enable seamless collaboration between team members
    - Include video conferencing and file sharing
    - Support offline work with automatic sync
    
    3. Create automated reporting dashboard
    - Real-time data visualization
    - Customizable metrics and KPIs
    - Automated alerts for critical issues
    
    Budget Considerations:
    Initial investment required for development team and infrastructure.
    Expected ROI within 12-18 months through efficiency gains.
    
    Timeline:
    Phase 1: Research and planning (2 months)
    Phase 2: Development and testing (6 months)  
    Phase 3: Launch and optimization (3 months)
    """
    
    # Create converter instance
    converter = OllamaDocumentToPowerPoint(ollama_model="qwen3:30b")
    
    # Create presentation from custom text
    print("Creating presentation from custom ideas...")
    output_file = converter.create_presentation_from_text(
        user_text=sample_text,
        presentation_title="Innovation Project Proposal",
        text_title="Project Ideas & Strategy",
        output_file="custom_ideas_presentation.pptx",
        enable_ai_analysis=True
    )
    
    if output_file:
        print(f"[SUCCESS] Custom text presentation created: {output_file}")
        print("[INFO] Features included:")
        print("• AI-enhanced analysis of your ideas")
        print("• Automatic content organization")
        print("• Executive summary generation") 
        print("• Actionable recommendations")
        print("• Professional slide formatting")
    else:
        print("[ERROR] Failed to create presentation from custom text")

def interactive_text_input_demo():
    """Interactive demo for text input functionality"""
    print("\n=== Interactive Text Input ===")
    print("Enter your ideas below (press Enter twice when finished):")
    print("You can include:")
    print("- Bullet points or numbered lists")
    print("- Multiple paragraphs") 
    print("- Project ideas, strategies, plans")
    print("- Any text content you want to present")
    print("\nYour input:")
    
    lines = []
    empty_lines = 0
    
    try:
        while True:
            line = input()
            if line.strip() == "":
                empty_lines += 1
                if empty_lines >= 2:  # Two empty lines = finished
                    break
            else:
                empty_lines = 0
                lines.append(line)
    except KeyboardInterrupt:
        print("\n[INFO] Input cancelled")
        return
    
    if not lines:
        print("[WARNING] No text entered")
        return
    
    user_text = "\n".join(lines)
    
    # Get presentation title
    try:
        title = input("\nEnter presentation title (or press Enter for default): ").strip()
        if not title:
            title = "My Ideas Presentation"
    except KeyboardInterrupt:
        title = "My Ideas Presentation"
    
    # Create presentation
    converter = OllamaDocumentToPowerPoint(ollama_model="qwen3:30b")
    
    print(f"\nCreating presentation: '{title}'...")
    output_file = converter.create_presentation_from_text(
        user_text=user_text,
        presentation_title=title,
        text_title="Custom Content",
        enable_ai_analysis=True
    )
    
    if output_file:
        print(f"[SUCCESS] Your presentation has been created: {output_file}")
    else:
        print("[ERROR] Failed to create presentation")

def create_document_summary():
    """
    DOCUMENT SUMMARIZER: Upload files and get structured summaries
    Perfect for creating presentations from existing documents
    """
    print("\n" + "📄" * 25)
    print("DOCUMENT SUMMARIZER - CREATE PRESENTATIONS FROM YOUR FILES")
    print("📄" * 25)
    print("\n📁 Drag and drop your files into this folder:")
    print("   - Word documents (.docx)")
    print("   - Excel files (.xlsx)")  
    print("   - PowerPoint files (.pptx)")
    print("   - PDF files (.pdf)")
    print("   - Text files (.txt)")
    print("\n⚡ No AI insights - Just clean document summaries!")
    
    # Get files from current directory
    import glob
    
    supported_files = []
    supported_files.extend(glob.glob("*.docx"))
    supported_files.extend(glob.glob("*.xlsx"))
    supported_files.extend(glob.glob("*.pptx"))
    supported_files.extend(glob.glob("*.pdf"))
    supported_files.extend(glob.glob("*.txt"))
    
    if not supported_files:
        print("\n❌ No supported files found in current directory")
        print("💡 Add some documents and try again!")
        return
    
    print(f"\n✅ Found {len(supported_files)} file(s):")
    for i, file in enumerate(supported_files, 1):
        print(f"  {i}. {file}")
    
    # Get presentation title
    title = input("\n📋 Presentation title: ").strip()
    if not title:
        title = "Document Summary Presentation"
    
    print(f"\n🚀 Creating summary presentation: '{title}'...")
    
    # Process files - with summarization enabled by default
    converter = OllamaDocumentToPowerPoint()
    
    # Process all files and create presentation
    all_content = converter.process_files(
        docx_files=[f for f in supported_files if f.endswith('.docx')],
        xlsx_files=[f for f in supported_files if f.endswith('.xlsx')],
        pptx_files=[f for f in supported_files if f.endswith('.pptx')],
        pdf_files=[f for f in supported_files if f.endswith('.pdf')],
        txt_files=[f for f in supported_files if f.endswith('.txt')]
    )
    
    if all_content:
        # Create presentation with document summaries
        output_file = converter.create_presentation(all_content, presentation_title=title)
        
        if output_file:
            print(f"\n🎉 SUCCESS! Created: {output_file}")
            print("📊 Your document summary presentation is ready!")
            print("\n📋 What's included:")
            print("  ✓ Document overviews")
            print("  ✓ Main points extracted") 
            print("  ✓ Key content summarized")
            print("  ✓ Action items identified")
        else:
            print("\n❌ Failed to create presentation")
    else:
        print("\n❌ No content could be extracted from files")

def easy_text_input():
    """
    EASIEST WAY: Simple text input interface
    Just run this function and type your ideas!
    """
    print("\n" + "🎯" * 25)
    print("EASY TEXT TO POWERPOINT CONVERTER")
    print("🎯" * 25)
    print("\n📝 Type your content below:")
    print("💡 Ideas, notes, plans, anything you want!")
    print("⏎ Press Enter twice when finished\n")
    
    # Collect text input
    lines = []
    empty_count = 0
    
    print("Start typing:")
    while empty_count < 2:
        try:
            line = input()
            if not line.strip():
                empty_count += 1
            else:
                empty_count = 0
            lines.append(line)
        except (KeyboardInterrupt, EOFError):
            break
    
    # Process the text
    user_text = "\n".join(lines[:-2]).strip()  # Remove last 2 empty lines
    
    if not user_text:
        print("❌ No text entered. Try again!")
        return
    
    print(f"\n✅ Got {len(user_text)} characters of content!")
    
    # Get title
    title = input("📋 Presentation title (or press Enter for default): ").strip()
    if not title:
        title = "My Ideas Presentation"
    
    # Create presentation
    print(f"\n🚀 Creating '{title}'...")
    
    converter = OllamaDocumentToPowerPoint()
    output_file = converter.create_presentation_from_text(
        user_text=user_text,
        presentation_title=title
    )
    
    if output_file:
        print(f"\n🎉 SUCCESS! Created: {output_file}")
        print("📂 Check your folder for the PowerPoint file!")
    else:
        print("\n❌ Failed to create presentation")

# Main execution block
if __name__ == "__main__":
    print("🚀 AI-Powered Document Summarization & PowerPoint Generator")
    print("=" * 65)
    print()
    print("Enhanced Features:")
    print("• Comprehensive document summarization with AI")
    print("• Executive summaries and detailed analysis")
    print("• [THEME] Random PowerPoint Design Themes:")
    print("  - Organic, Ion, Retrospect, Savon, Slice, Wisp")
    print("  - Berlin, Celestial, Facet, Gallery, Headlines")
    print("  - Integral, Metropolitan, Office Theme")
    print("• [CHART] Perfect Chart Fitting - No overlapping content!")
    print("• 📐 Smart text sizing - Everything fits perfectly")
    print("• Extract from Word (.docx) and OpenOffice (.odt) documents")
    print("• Advanced Excel analysis with charts (.xlsx)")
    print("• Power BI file processing (.pbix)")
    print("• Ollama AI with Qwen3:30B for intelligent insights")
    print("• Smart content fitting - no overflow text!")
    print("• Local AI - no external API keys required!")
    print()
    print("Setup Requirements:")
    print("1. Ollama installed and running on your system")
    print("2. Qwen3:30B model available in Ollama")
    print("3. Files to process (already found in wrapper test folder)")
    print()
    print("Available methods:")
    print("1. converter.process_files(input_folder='path/to/folder')")
    print("2. converter.process_files(word_files=[...], excel_files=[...], odt_files=[...])")
    print()
    print("🎯 Ready to process your files with Ollama AI!")
    print()
    print("Choose an option:")
    print("1. Process files from current folder (auto-detect)")
    print("2. Easy text input (type your own content)")
    print("3. Interactive file selection")
    print("4. Exit")
    print()
    
    try:
        choice = input("Enter your choice (1-4): ").strip()
        
        if choice == "1":
            print("\n🚀 Auto-processing files in current directory...")
            wrapper_test_folder = r"c:\Users\aidan\OneDrive\Desktop\AI POWERPOINT WRAPPER\wrapper test"
            
            converter = OllamaDocumentToPowerPoint(ollama_model="qwen3:30b")
            print(f"[SEARCH] Scanning folder: {wrapper_test_folder}")
            all_content = converter.process_files(input_folder=wrapper_test_folder)
            
            if all_content:
                print(f"[FOLDER] Found {len(all_content)} file(s) to process")
                output_file = converter.create_presentation(
                    all_content, 
                    "powerpoint_project.pptx",
                    "PowerPoint Project - Document Analysis"
                )
                print(f"[OK] Complete presentation created: {output_file}")
            else:
                print("[ERROR] No files found to process")
                
        elif choice == "2":
            easy_text_input()
            
        elif choice == "3":
            create_document_summary()
            
        elif choice == "4":
            print("👋 Goodbye!")
            
        else:
            print("❌ Invalid choice. Please run the script again.")
            
    except (KeyboardInterrupt, EOFError):
        print("\n👋 Goodbye!")